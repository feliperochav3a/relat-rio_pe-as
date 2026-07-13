"""
PPT Builder — gera Relatório de Peças com base no template V3A.
Lógica determinística, sem LLM.
"""
import copy
import io
import re
from pathlib import Path

from PIL import Image as PILImage
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

TEMPLATE_PATH = Path(__file__).parent.parent / "templates" / "template.pptx"

# Índices dos slides no template original
IDX_COVER    = 0
IDX_PORTRAIT = 1   # Sacochila Amarela — layout retrato
IDX_LANDSCAPE = 3  # Bike              — layout paisagem
IDX_TOTAL    = 4
IDX_LAST     = 5

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


# ── Helpers de texto ─────────────────────────────────────────────────────────

def filename_to_title(filename: str) -> str:
    """'barraca_verde.jpg' → 'BARRACA VERDE'"""
    stem = Path(filename).stem
    title = re.sub(r"[_\-]+", " ", stem)
    title = re.sub(r"\s+", " ", title).strip()
    return title.upper()


def format_subtitle(subtitle: str) -> str:
    """
    Garante o padrão 'CLIENTE | CAMPANHA'.
    Se o usuário não incluiu '|', tenta detectar a fronteira.
    """
    subtitle = subtitle.strip().upper()
    if "|" in subtitle:
        return subtitle
    # Heurística: palavras que sugerem local/cliente vêm antes
    # Divide após o 2º token se houver 3+ tokens
    parts = subtitle.split()
    if len(parts) >= 3:
        client = " ".join(parts[:2])
        campaign = " ".join(parts[2:])
        return f"{client} | {campaign}"
    return subtitle


# ── Helpers de imagem ────────────────────────────────────────────────────────

def image_orientation(path: str) -> str:
    """Retorna 'portrait' ou 'landscape' baseado nas dimensões."""
    with PILImage.open(path) as img:
        w, h = img.size
    return "landscape" if w > h else "portrait"


def _blip_of(pic_shape):
    """Encontra o elemento blip de um shape de imagem."""
    return pic_shape.element.find(f".//{{{NS_A}}}blip")


def replace_image(slide, pic_shape, new_image_path: str):
    """
    Substitui a imagem mantendo a proporção original (fit/contain no placeholder).
    A imagem é centralizada dentro da área reservada pelo template.
    """
    blip = _blip_of(pic_shape)
    if blip is None:
        return

    old_rid = blip.get(f"{{{NS_R}}}embed")

    # Dimensões reais da imagem
    with PILImage.open(new_image_path) as img:
        img_w, img_h = img.size

    # Área do placeholder no template (em EMU)
    ph_left = pic_shape.left
    ph_top  = pic_shape.top
    ph_w    = pic_shape.width
    ph_h    = pic_shape.height

    # Escala proporcional (contain: imagem inteira visível, sem distorção)
    scale = min(ph_w / img_w, ph_h / img_h)
    fit_w = int(img_w * scale)
    fit_h = int(img_h * scale)

    # Centraliza dentro do placeholder
    left = ph_left + (ph_w - fit_w) // 2
    top  = ph_top  + (ph_h - fit_h) // 2

    # Adiciona imagem temporária no tamanho correto para obter o ImagePart/rId
    tmp = slide.shapes.add_picture(new_image_path, left, top, fit_w, fit_h)
    new_rid = _blip_of(tmp).get(f"{{{NS_R}}}embed")

    # Redireciona o blip original para a nova imagem
    blip.set(f"{{{NS_R}}}embed", new_rid)

    # Atualiza posição e tamanho do shape no XML para refletir o fit
    xfrm = pic_shape.element.find(f".//{{{NS_A}}}xfrm")
    if xfrm is not None:
        off = xfrm.find(f"{{{NS_A}}}off")
        ext = xfrm.find(f"{{{NS_A}}}ext")
        if off is not None:
            off.set("x", str(left))
            off.set("y", str(top))
        if ext is not None:
            ext.set("cx", str(fit_w))
            ext.set("cy", str(fit_h))

    # Remove o shape temporário (ImagePart permanece no pacote)
    slide.shapes._spTree.remove(tmp.element)

    # Descarta a relação antiga
    if old_rid and old_rid != new_rid:
        try:
            slide.part.drop_rel(old_rid)
        except Exception:
            pass


# ── Clone / delete / reorder de slides ──────────────────────────────────────

def clone_slide(prs: Presentation, source_index: int):
    """
    Clona o slide em source_index e adiciona ao final da apresentação.
    Copia relações de imagens e substitui rIds no XML copiado.
    """
    src = prs.slides[source_index]

    # Adiciona slide com o mesmo layout
    new_slide = prs.slides.add_slide(src.slide_layout)

    # Copia relações de imagem → mapeia rId antigo para novo
    rid_map: dict[str, str] = {}
    for rId, rel in src.part.rels.items():
        if "image" in rel.reltype:
            new_rid = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rid_map[rId] = new_rid

    # Deepcopy do spTree e substitui rIds
    src_tree = copy.deepcopy(src.shapes._spTree)
    _substitute_rids(src_tree, rid_map)

    # Substitui conteúdo do novo slide
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree)[2:]:      # mantém nvGrpSpPr e grpSpPr
        sp_tree.remove(child)
    for child in list(src_tree)[2:]:
        sp_tree.append(child)

    return new_slide


def _substitute_rids(element, rid_map: dict):
    """Substitui rIds em todos os atributos do elemento (recursivo)."""
    for el in element.iter():
        for attr in (f"{{{NS_R}}}embed", f"{{{NS_R}}}link", f"{{{NS_R}}}id"):
            val = el.get(attr)
            if val and val in rid_map:
                el.set(attr, rid_map[val])


def delete_slide(prs: Presentation, index: int):
    """Remove o slide no índice informado."""
    sldIdLst = prs.slides._sldIdLst
    slide_el = sldIdLst[index]
    rid = slide_el.get(f"{{{NS_R}}}id")
    sldIdLst.remove(slide_el)
    try:
        prs.part.drop_rel(rid)
    except Exception:
        pass


def reorder_slides(prs: Presentation, new_order: list[int]):
    """Reordena slides conforme lista de índices."""
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    for item in items:
        sldIdLst.remove(item)
    for i in new_order:
        sldIdLst.append(items[i])


# ── Manipulação de conteúdo dos slides ──────────────────────────────────────

def _find_shape(shapes, name: str):
    """Busca shape pelo nome (recursivo em grupos)."""
    for s in shapes:
        if s.name == name:
            return s
        if s.shape_type == 6:  # GROUP
            found = _find_shape(s.shapes, name)
            if found:
                return found
    return None


def set_cover_subtitle(slide, subtitle: str):
    """Atualiza o subtítulo da capa ('CaixaDeTexto 14' dentro do grupo 'Agrupar 2')."""
    group = _find_shape(slide.shapes, "Agrupar 2")
    if not group:
        return
    shape = _find_shape(group.shapes, "CaixaDeTexto 14")
    if not shape or not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.text = subtitle
        if para.runs:
            break   # só a primeira linha não-vazia


def set_piece_title(slide, title: str):
    """
    Define o título da peça ('CaixaDeTexto 11').
    Divide em até 2 linhas automaticamente.
    """
    shape = _find_shape(slide.shapes, "CaixaDeTexto 11")
    if not shape or not shape.has_text_frame:
        return

    words = title.upper().split()
    if len(words) <= 1:
        lines = words or [""]
    elif len(words) == 2:
        lines = words
    else:
        mid = (len(words) + 1) // 2
        lines = [" ".join(words[:mid]), " ".join(words[mid:])]

    txBody = shape.text_frame._txBody
    existing = txBody.findall(f"{{{NS_A}}}p")

    # Template de parágrafo para copiar a formatação
    tmpl = copy.deepcopy(existing[0]) if existing else None

    # Remove parágrafos antigos
    for p in existing:
        txBody.remove(p)

    for line in lines:
        if tmpl is not None:
            new_p = copy.deepcopy(tmpl)
        else:
            from lxml import etree
            new_p = etree.Element(f"{{{NS_A}}}p")

        # Limpa runs do parágrafo copiado
        for r in new_p.findall(f"{{{NS_A}}}r"):
            new_p.remove(r)

        # Cria run com o texto da linha copiando a formatação do template
        if tmpl is not None:
            tmpl_runs = tmpl.findall(f"{{{NS_A}}}r")
            if tmpl_runs:
                new_r = copy.deepcopy(tmpl_runs[0])
                t_el = new_r.find(f"{{{NS_A}}}t")
                if t_el is None:
                    from lxml import etree
                    t_el = etree.SubElement(new_r, f"{{{NS_A}}}t")
                t_el.text = line
                new_p.append(new_r)
        else:
            from lxml import etree
            new_r = etree.SubElement(new_p, f"{{{NS_A}}}r")
            t_el = etree.SubElement(new_r, f"{{{NS_A}}}t")
            t_el.text = line

        txBody.append(new_p)


def set_total(slide, n: int):
    """Atualiza o texto do slide de total ('CaixaDeTexto 11')."""
    shape = _find_shape(slide.shapes, "CaixaDeTexto 11")
    if not shape or not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if "total" in run.text.lower() or "peça" in run.text.lower():
                run.text = f"Total de {n:02d} peças"
                return
        # Fallback: atualiza qualquer run não-vazio
        for run in para.runs:
            if run.text.strip():
                run.text = f"Total de {n:02d} peças"
                return


def _main_picture(slide):
    """
    Retorna o shape de peça do slide.
    Busca primeiro pelo nome do template ('Imagem 12' retrato, 'Imagem 5' paisagem),
    depois pela maior área como fallback.
    """
    pics = [s for s in slide.shapes if s.shape_type == 13]
    if not pics:
        return None
    for name in ("Imagem 12", "Imagem 5", "Imagem 2"):
        match = next((s for s in pics if s.name == name), None)
        if match:
            return match
    # Fallback: maior área excluindo shapes muito pequenos (logos, ícones < 1")
    from pptx.util import Inches
    large = [s for s in pics if s.width > Inches(1) and s.height > Inches(1)]
    if large:
        return max(large, key=lambda s: s.width * s.height)
    return max(pics, key=lambda s: s.width * s.height)


# ── Função principal ─────────────────────────────────────────────────────────

def build_report(subtitle: str, pieces: list[dict]) -> bytes:
    """
    Gera o PPTX e retorna os bytes.

    pieces: [{"filename": "barraca.png", "image_path": "/tmp/barraca.png"}, ...]
    """
    prs = Presentation(str(TEMPLATE_PATH))
    n = len(pieces)

    # 1. Subtítulo da capa
    set_cover_subtitle(prs.slides[IDX_COVER], format_subtitle(subtitle))

    # 2. Clona slides de peça para cada peça do usuário (appended ao final)
    for piece in pieces:
        orientation = image_orientation(piece["image_path"])
        src_idx = IDX_PORTRAIT if orientation == "portrait" else IDX_LANDSCAPE
        new_slide = clone_slide(prs, src_idx)
        title = filename_to_title(piece["filename"])
        set_piece_title(new_slide, title)
        pic = _main_picture(new_slide)
        if pic:
            replace_image(new_slide, pic, piece["image_path"])

    # 3. Deleta os 3 slides de peça originais (do maior ao menor para não deslocar índices)
    for idx in (3, 2, 1):
        delete_slide(prs, idx)

    # Após deleção: [0=cover, 1=total, 2=last, 3..2+n=novas_peças]
    # Ordem desejada: [cover, peça1..peçaN, total, last]
    new_order = [0] + list(range(3, 3 + n)) + [1, 2]
    reorder_slides(prs, new_order)

    # 4. Atualiza total (agora em índice n+1)
    set_total(prs.slides[n + 1], n)

    # 5. Salva em memória
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ── Diagnóstico para o Diretor de Arte ───────────────────────────────────────

def inspect_report(pptx_bytes: bytes) -> dict:
    """
    Lê o PPTX gerado e retorna um relatório estrutural para o Diretor de Arte.
    """
    prs = Presentation(io.BytesIO(pptx_bytes))
    slides_info = []

    for i, slide in enumerate(prs.slides):
        texts = []
        has_image = False
        fonts = []

        def walk(shapes):
            nonlocal has_image
            for s in shapes:
                if s.shape_type == 13:
                    has_image = True
                if s.has_text_frame:
                    for para in s.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                        for run in para.runs:
                            if run.font.name:
                                fonts.append(run.font.name)
                if s.shape_type == 6:
                    walk(s.shapes)

        walk(slide.shapes)
        slides_info.append({
            "index": i,
            "texts": texts,
            "has_image": has_image,
            "fonts": list(set(fonts)),
        })

    return {
        "total_slides": len(prs.slides),
        "slides": slides_info,
    }
